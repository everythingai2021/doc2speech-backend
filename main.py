from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.responses import JSONResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pyt2s.services import stream_elements
from pptx import Presentation
import os
import io
import string
import PyPDF2

import tempfile

app = FastAPI()

origins = [
    "https://doc2speech.vercel.app",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins, 
    allow_credentials=True,
    allow_methods=["*"],  
    allow_headers=["*"],  
)

@app.get("/")
def root():
    return {"Hello": "World"}

@app.post("/convert")
async def convert(file: UploadFile = File(...)):
    if not (file.filename.endswith(".pptx") or file.filename.endswith(".pdf")):
        raise HTTPException(status_code=400, detail="File must be a .pptx PowerPoint file or a .pdf file.")

    print(f"Received file: {file.filename}")

    contents = await file.read()
    if file.filename.endswith(".pdf"):
        return await parse_pdf(contents)
    elif file.filename.endswith(".pptx"):
        return await parse_pptx(contents)


async def parse_pdf(contents):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
            pdf_path = temp_pdf.name
            temp_pdf.write(contents)

        extracted_text = extract_text(pdf_path)
        print(f"Extracted text: {extracted_text[:100]}...")
        if len(extracted_text.strip()) == 0:
            return JSONResponse(content={"error": "No extractable text found in the PDF."})

        data = stream_elements.requestTTS(extracted_text, stream_elements.Voice.Amy.value)

        return StreamingResponse(
            io.BytesIO(data),
            media_type="audio/mpeg",
            headers={"Content-Disposition": "inline; filename=output.mp3"}
        )

    except Exception as e:
        print(f"Error processing PDF: {str(e)}")
        return JSONResponse(content={"error": f"An error occurred while processing the PDF: {str(e)}"})

    finally:
        if os.path.exists(pdf_path):
            os.remove(pdf_path)

def extract_pptx_text(pptx_path):
    text = ""
    ppt = Presentation(pptx_path)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                text += clean_text(shape.text_frame.text) + "\n"
    return text

async def parse_pptx(contents):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_pptx:
            pptx_path = temp_pptx.name
            temp_pptx.write(contents)

        extracted_text = extract_pptx_text(pptx_path)
        print(f"Extracted text: {extracted_text[:100]}...")
        
        if len(extracted_text.strip()) == 0:
            return JSONResponse(content={"error": "No extractable text found in the PowerPoint."})

        data = stream_elements.requestTTS(extracted_text, stream_elements.Voice.Amy.value)

        return StreamingResponse(
            io.BytesIO(data),
            media_type="audio/mpeg",
            headers={"Content-Disposition": "inline; filename=output.mp3"}
        )
    except Exception as e:
        print(f"Error processing PPTX: {str(e)}")
        return JSONResponse(content={"error": f"An error occurred while processing the PPTX: {str(e)}"})

    finally:
        if os.path.exists(pptx_path):
            os.remove(pptx_path)

def clean_text(text):
    allowed = string.ascii_letters + string.digits + string.punctuation + " \n"
    strings = [[]]
    idx = 0
    for c in text:
        if c in allowed:
            if c == "\n":
                strings[idx].append(".")
                strings[idx].append(c)
                idx += 1
                strings.append([])
            else:
                strings[idx].append(c)
    res = ""
    for s in strings:
        res += "".join(s)
    return res

def extract_text(pdf_path):
    text = ""
    with open(pdf_path, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text += clean_text(page_text) + "\n"
            else:
                text += "[No text found on this page]\n"
    return text


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)

